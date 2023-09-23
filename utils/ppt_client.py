from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor as RGB

def create_slide(prs, title, content, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.shapes.placeholders[1]
    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.text = content


def create_image_slide(prs, img_path, layout_index=6):
    slide_layout = prs.slide_layouts[layout_index]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(5.5))

def create_text_slide(prs, text, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    body_placeholder = slide.shapes.placeholders[1]
    tf = body_placeholder.text_frame
    tf.text = text

def generate_catalogue(catalogue):
    prs = Presentation()
    for item in catalogue:
        img_path = f"images/{item['cod_sap']}.jpg"  # Assuming images are in .jpg format
        create_image_slide(prs, img_path)
        create_text_slide(prs, item['sales_pitch'])

    prs.save("catalogue.pptx")


def main():
    prs = Presentation()

    # Slide 1: Title Slide
    # slide_layout = prs.slide_layouts[0]  # Title layout
    # slide = prs.slides.add_slide(slide_layout)
    # title = slide.shapes.title
    # subtitle = slide.placeholders[1]
    # title.text = "First-Time Homebuyers Guide"
    # subtitle.text = "Understanding the Homebuying Process in Canada"

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Title layout
    slide = prs.slides.add_slide(slide_layout)
    # slide.background.fill.solid()
    # slide.background.fill.fore_color.rgb = RGB(255, 223, 186)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "First-Time Homebuyers Guide"
    subtitle.text = "Understanding the Homebuying Process in Canada"
    slide.shapes.add_picture('placeholder.jpg', Inches(2), Inches(2), height=Inches(5.5))


    # Slide 2: Introduction
    create_slide(prs, "Introduction", 
                 "Buying a home for the first time can be both exciting and overwhelming. This guide will help you navigate the process in Canada.")

    # Slide 3: Determine Your Budget
    create_slide(prs, "Determine Your Budget", 
                 "Before starting your home search, determine how much you can afford. Consider getting a mortgage pre-approval to understand your budget.")

    # Slide 4: Choose a Realtor
    create_slide(prs, "Choose a Realtor", 
                 "A realtor can provide valuable insights, help you find properties, and negotiate on your behalf.")

    # Slide 5: Home Inspection
    create_slide(prs, "Home Inspection", 
                 "Once you've chosen a home, get it inspected. This will identify any potential issues or repairs needed.")

    # Slide 6: Understand the Costs
    create_slide(prs, "Understand the Costs", 
                 "Apart from the home price, consider other costs like land transfer tax, legal fees, and home insurance.")

    # Slide 7: Mortgage and Financing
    create_slide(prs, "Mortgage and Financing", 
                 "Understand the different types of mortgages available, interest rates, and choose the best fit for your financial situation.")

    # Slide 8: Closing the Deal
    create_slide(prs, "Closing the Deal", 
                 "Once everything is in place, you'll finalize the deal. Ensure all paperwork is in order and understand your responsibilities.")

    # Save the presentation
    prs.save('First_Time_Homebuyers_Guide.pptx')

if __name__ == "__main__":
    # main()
    catalogue = [{'cod_sap': '200098377', 'sales_pitch': '...Description...'}, {'cod_sap': '200105655', 'sales_pitch':'item2'}]  # Your catalogue items
    generate_catalogue(catalogue)